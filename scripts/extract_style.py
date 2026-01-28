#!/usr/bin/env python3
"""
Extract visual style (colors, fonts, layouts) from a PowerPoint presentation.
Outputs a JSON style definition that can be used to generate matching CSS.

Usage:
    python extract_style.py input.pptx [output.json]

Requires: pip install python-pptx Pillow
"""

import sys
import json
import os
from pathlib import Path
from collections import Counter
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
    import io
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


def rgb_to_hex(rgb_color) -> Optional[str]:
    """Convert pptx RGBColor to hex string."""
    if rgb_color is None:
        return None
    if isinstance(rgb_color, RGBColor):
        return f"#{rgb_color.red:02x}{rgb_color.green:02x}{rgb_color.blue:02x}"
    return None


def get_color_from_fill(fill) -> Optional[str]:
    """Extract color from a fill object."""
    try:
        if fill is None:
            return None
        if fill.type is not None:
            # Solid fill
            if hasattr(fill, 'fore_color') and fill.fore_color:
                if fill.fore_color.type == 1:  # RGB
                    return rgb_to_hex(fill.fore_color.rgb)
                elif fill.fore_color.theme_color is not None:
                    # Theme color - return a marker
                    return f"theme:{fill.fore_color.theme_color}"
    except Exception:
        pass
    return None


def get_font_info(font) -> dict:
    """Extract font information."""
    info = {}
    try:
        if font.name:
            info['name'] = font.name
        if font.size:
            info['size_pt'] = font.size.pt
        if font.bold:
            info['bold'] = True
        if font.italic:
            info['italic'] = True
        if font.color and font.color.rgb:
            info['color'] = rgb_to_hex(font.color.rgb)
    except Exception:
        pass
    return info


def extract_theme_colors(prs: Presentation) -> dict:
    """Extract theme colors from the presentation."""
    theme_colors = {}
    try:
        # Access the theme through the slide master
        if prs.slide_masters:
            master = prs.slide_masters[0]
            # Theme colors are in the theme part
            theme = master.part.slide_master.get_or_add_cSld()
            # Note: Direct theme color access is limited in python-pptx
            # We'll extract colors from actual content instead
    except Exception:
        pass
    return theme_colors


def extract_background_color(slide) -> Optional[str]:
    """Extract background color from a slide."""
    try:
        background = slide.background
        if background.fill:
            return get_color_from_fill(background.fill)
    except Exception:
        pass
    return None


def extract_dominant_image_colors(prs: Presentation, max_images: int = 5) -> list:
    """Extract dominant colors from images in the presentation."""
    if not HAS_PIL:
        return []

    colors = []
    image_count = 0

    for slide in prs.slides:
        if image_count >= max_images:
            break
        for shape in slide.shapes:
            if image_count >= max_images:
                break
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_bytes = shape.image.blob
                    img = Image.open(io.BytesIO(image_bytes))
                    img = img.convert('RGB')
                    img = img.resize((50, 50))  # Downsample for speed

                    # Get most common colors
                    pixels = list(img.getdata())
                    color_counts = Counter(pixels)
                    for color, count in color_counts.most_common(3):
                        hex_color = f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"
                        colors.append(hex_color)
                    image_count += 1
                except Exception:
                    pass

    return colors


def analyze_presentation(pptx_path: str) -> dict:
    """Analyze a PowerPoint file and extract style information."""
    prs = Presentation(pptx_path)

    # Collect data
    background_colors = []
    text_colors = []
    fonts = []
    font_sizes = []
    layouts_used = []

    for slide in prs.slides:
        # Background
        bg_color = extract_background_color(slide)
        if bg_color and not bg_color.startswith('theme:'):
            background_colors.append(bg_color)

        # Layout name
        if slide.slide_layout:
            layouts_used.append(slide.slide_layout.name)

        # Shapes
        for shape in slide.shapes:
            # Text content
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_info = get_font_info(run.font)
                        if 'name' in font_info:
                            fonts.append(font_info['name'])
                        if 'size_pt' in font_info:
                            font_sizes.append(font_info['size_pt'])
                        if 'color' in font_info:
                            text_colors.append(font_info['color'])

            # Shape fills (for accent colors)
            if hasattr(shape, 'fill'):
                fill_color = get_color_from_fill(shape.fill)
                if fill_color and not fill_color.startswith('theme:'):
                    text_colors.append(fill_color)

    # Analyze collected data
    bg_counts = Counter(background_colors)
    text_counts = Counter(text_colors)
    font_counts = Counter(fonts)
    size_counts = Counter(font_sizes)
    layout_counts = Counter(layouts_used)

    # Determine primary colors
    primary_bg = bg_counts.most_common(1)[0][0] if bg_counts else "#ffffff"

    # Categorize text colors by frequency and darkness
    sorted_text_colors = text_counts.most_common(10)

    # Extract image colors for accent detection
    image_colors = extract_dominant_image_colors(prs)

    # Determine if dark or light theme
    def color_brightness(hex_color):
        hex_color = hex_color.lstrip('#')
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        return (r * 299 + g * 587 + b * 114) / 1000

    bg_brightness = color_brightness(primary_bg)
    is_dark_theme = bg_brightness < 128

    # Separate text colors into primary, secondary, accent
    light_colors = []
    dark_colors = []
    mid_colors = []

    for color, count in sorted_text_colors:
        brightness = color_brightness(color)
        if brightness > 200:
            light_colors.append((color, count))
        elif brightness < 55:
            dark_colors.append((color, count))
        else:
            mid_colors.append((color, count))

    # Assign colors based on theme
    if is_dark_theme:
        text_primary = light_colors[0][0] if light_colors else "#ffffff"
        text_secondary = mid_colors[0][0] if mid_colors else "#9ca3af"
    else:
        text_primary = dark_colors[0][0] if dark_colors else "#1f2937"
        text_secondary = mid_colors[0][0] if mid_colors else "#6b7280"

    # Find accent color (most saturated non-grayscale color)
    def color_saturation(hex_color):
        hex_color = hex_color.lstrip('#')
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        max_c = max(r, g, b)
        min_c = min(r, g, b)
        if max_c == 0:
            return 0
        return (max_c - min_c) / max_c

    accent_candidates = [(c, color_saturation(c)) for c, _ in sorted_text_colors if color_saturation(c) > 0.3]
    accent_candidates.extend([(c, color_saturation(c)) for c in image_colors if color_saturation(c) > 0.3])
    accent_candidates.sort(key=lambda x: x[1], reverse=True)
    accent_color = accent_candidates[0][0] if accent_candidates else "#3b82f6"

    # Determine fonts
    font_priority = font_counts.most_common(3)
    display_font = font_priority[0][0] if font_priority else "Inter"
    body_font = font_priority[1][0] if len(font_priority) > 1 else display_font

    # Determine typical font sizes
    size_priority = size_counts.most_common(5)
    heading_sizes = [s for s, c in size_priority if s >= 24]
    body_sizes = [s for s, c in size_priority if 12 <= s < 24]

    heading_size = max(heading_sizes) if heading_sizes else 48
    body_size = max(body_sizes) if body_sizes else 18

    # Build style definition
    style = {
        "metadata": {
            "source_file": os.path.basename(pptx_path),
            "slide_count": len(prs.slides),
            "theme_type": "dark" if is_dark_theme else "light",
            "layouts_detected": list(layout_counts.keys())
        },
        "colors": {
            "bg_primary": primary_bg,
            "bg_secondary": adjust_color(primary_bg, 10 if is_dark_theme else -10),
            "text_primary": text_primary,
            "text_secondary": text_secondary,
            "accent": accent_color,
            "accent_glow": accent_color + "40"  # 25% opacity for glow
        },
        "typography": {
            "font_display": display_font,
            "font_body": body_font,
            "heading_size_pt": heading_size,
            "body_size_pt": body_size,
            "google_fonts_url": generate_google_fonts_url([display_font, body_font])
        },
        "spacing": {
            "slide_padding": "clamp(2rem, 5vw, 4rem)"
        },
        "animation": {
            "duration_normal": "0.6s",
            "ease_function": "cubic-bezier(0.16, 1, 0.3, 1)"
        },
        "raw_data": {
            "all_background_colors": list(bg_counts.keys())[:5],
            "all_text_colors": list(text_counts.keys())[:10],
            "all_fonts": list(font_counts.keys()),
            "all_font_sizes": sorted(list(size_counts.keys()), reverse=True)[:10]
        }
    }

    return style


def adjust_color(hex_color: str, amount: int) -> str:
    """Lighten (positive) or darken (negative) a hex color."""
    hex_color = hex_color.lstrip('#')
    r = max(0, min(255, int(hex_color[0:2], 16) + amount))
    g = max(0, min(255, int(hex_color[2:4], 16) + amount))
    b = max(0, min(255, int(hex_color[4:6], 16) + amount))
    return f"#{r:02x}{g:02x}{b:02x}"


def generate_google_fonts_url(fonts: list) -> str:
    """Generate a Google Fonts URL for the given fonts."""
    # Common font mappings (fonts that need special handling)
    google_font_names = {
        'Arial': None,  # System font
        'Helvetica': None,  # System font
        'Times New Roman': None,  # System font
        'Calibri': 'Open Sans',  # Close alternative
        'Cambria': 'Merriweather',  # Close alternative
    }

    font_params = []
    for font in fonts:
        if font in google_font_names:
            if google_font_names[font]:
                font_params.append(google_font_names[font].replace(' ', '+') + ':wght@400;600;700')
        else:
            font_params.append(font.replace(' ', '+') + ':wght@400;600;700')

    if not font_params:
        return "https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700&display=swap"

    return f"https://fonts.googleapis.com/css2?family={'&family='.join(font_params)}&display=swap"


def generate_css_variables(style: dict) -> str:
    """Generate CSS custom properties from style definition."""
    colors = style['colors']
    typography = style['typography']
    spacing = style['spacing']
    animation = style['animation']

    css = f""":root {{
    /* Colors - extracted from {style['metadata']['source_file']} */
    --bg-primary: {colors['bg_primary']};
    --bg-secondary: {colors['bg_secondary']};
    --text-primary: {colors['text_primary']};
    --text-secondary: {colors['text_secondary']};
    --accent: {colors['accent']};
    --accent-glow: {colors['accent_glow']};

    /* Typography */
    --font-display: '{typography['font_display']}', sans-serif;
    --font-body: '{typography['font_body']}', sans-serif;

    /* Spacing */
    --slide-padding: {spacing['slide_padding']};

    /* Animation */
    --duration-normal: {animation['duration_normal']};
    --ease-out-expo: {animation['ease_function']};
}}"""
    return css


def generate_image_style_json(style: dict) -> dict:
    """Generate image-style.json content for image generation."""
    colors = style['colors']
    theme_type = style['metadata']['theme_type']

    mood = "professional/clean" if theme_type == "light" else "modern/dramatic"

    return {
        "signature": f"minimalist illustration, {theme_type} theme, colors ({colors['bg_primary']}, {colors['accent']}), soft shadows, editorial feel",
        "mood": mood,
        "color_palette": {
            "primary": colors['text_primary'],
            "secondary": colors['text_secondary'],
            "accent": colors['accent'],
            "background": colors['bg_primary']
        },
        "style_keywords": ["minimalist", "illustration", "soft shadows", "editorial", theme_type],
        "negative_prompt": "photorealistic, 3D render, cluttered, neon, stock photo"
    }


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_style.py input.pptx [output.json]")
        print("\nExtracts visual style (colors, fonts) from a PowerPoint presentation.")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    print(f"Analyzing: {pptx_path}")
    style = analyze_presentation(pptx_path)

    # Output
    if output_path:
        with open(output_path, 'w') as f:
            json.dump(style, f, indent=2)
        print(f"Style definition written to: {output_path}")
    else:
        print("\n" + "="*60)
        print("EXTRACTED STYLE DEFINITION")
        print("="*60)
        print(json.dumps(style, indent=2))

    # Always print CSS variables
    print("\n" + "="*60)
    print("CSS CUSTOM PROPERTIES")
    print("="*60)
    print(generate_css_variables(style))

    # Print image style
    print("\n" + "="*60)
    print("IMAGE STYLE (for .claude-design/image-style.json)")
    print("="*60)
    print(json.dumps(generate_image_style_json(style), indent=2))

    # Print Google Fonts link
    print("\n" + "="*60)
    print("GOOGLE FONTS LINK")
    print("="*60)
    print(f'<link rel="stylesheet" href="{style["typography"]["google_fonts_url"]}">')


if __name__ == "__main__":
    main()
