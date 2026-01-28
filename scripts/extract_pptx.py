#!/usr/bin/env python3
"""
extract_pptx.py - Extract content and images from PowerPoint files
"""

import json
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_pptx(file_path: str, output_dir: str) -> dict:
    """
    Extract all content from a PowerPoint file.

    Returns a dict with slides, text, images, and speaker notes.
    """
    prs = Presentation(file_path)
    output_path = Path(output_dir)
    assets_dir = output_path / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    slides_data = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "number": slide_num,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
            "layout": get_layout_name(slide)
        }

        # Extract shapes
        for shape in slide.shapes:
            # Title
            if shape.is_placeholder and shape == slide.shapes.title:
                if shape.has_text_frame:
                    slide_data["title"] = shape.text.strip()

            # Text content
            elif shape.has_text_frame:
                text = extract_text_frame(shape.text_frame)
                if text:
                    slide_data["content"].append({
                        "type": "text",
                        "content": text
                    })

            # Tables
            elif shape.has_table:
                table_data = extract_table(shape.table)
                slide_data["content"].append({
                    "type": "table",
                    "content": table_data
                })

            # Images
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_info = extract_image(shape, slide_num, len(slide_data["images"]) + 1, assets_dir)
                if image_info:
                    slide_data["images"].append(image_info)

        # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                slide_data["notes"] = notes_frame.text.strip()

        slides_data.append(slide_data)

    result = {
        "source_file": os.path.basename(file_path),
        "slide_count": len(slides_data),
        "slides": slides_data
    }

    # Save extraction result
    json_path = output_path / "extracted_content.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=2, ensure_ascii=False)

    print(f"Extracted {len(slides_data)} slides to {output_dir}")
    print(f"Content saved to {json_path}")

    return result


def get_layout_name(slide) -> str:
    """Get the layout name for slide type hints."""
    try:
        return slide.slide_layout.name
    except:
        return "unknown"


def extract_text_frame(text_frame) -> list:
    """Extract paragraphs from a text frame."""
    paragraphs = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append({
                "text": text,
                "level": para.level  # Indentation level (0 = top level)
            })
    return paragraphs


def extract_table(table) -> dict:
    """Extract table data as rows and columns."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip())
        rows.append(cells)

    return {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "data": rows
    }


def extract_image(shape, slide_num: int, img_num: int, assets_dir: Path) -> dict:
    """Extract and save an image from a shape."""
    try:
        image = shape.image
        ext = image.ext
        filename = f"slide{slide_num:02d}_img{img_num}.{ext}"
        filepath = assets_dir / filename

        with open(filepath, "wb") as f:
            f.write(image.blob)

        return {
            "path": f"assets/{filename}",
            "width": shape.width,
            "height": shape.height,
            "alt": f"Image from slide {slide_num}"
        }
    except Exception as e:
        print(f"Warning: Could not extract image from slide {slide_num}: {e}")
        return None


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <input.pptx> [output_dir]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "."

    extract_pptx(input_file, output_dir)
